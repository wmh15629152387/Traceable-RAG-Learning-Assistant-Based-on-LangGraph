from __future__ import annotations
import io
from typing import Any


class DocumentLoader:
    """
    Load PPT/PDF/Markdown/Text into unified doc dicts.
    For demo: use basic extraction; you can swap to unstructured / pymupdf / python-pptx later.
    """

    def load_bytes(self, filename: str, data: bytes, content_type: str) -> list[dict[str, Any]]:
        name = filename.lower()

        if name.endswith(".md") or content_type in {"text/markdown", "text/plain"}:
            text = data.decode("utf-8", errors="ignore")
            return [{"doc_title": filename, "text": text, "locator": None, "source_type": "markdown"}]

        if name.endswith(".pdf") or content_type == "application/pdf":
            # Minimal PDF extraction (optional dependency)
            try:
                import pypdf  # type: ignore
                reader = pypdf.PdfReader(io.BytesIO(data))
                docs = []
                for i, page in enumerate(reader.pages):
                    text = page.extract_text() or ""
                    docs.append({"doc_title": filename, "text": text, "locator": f"page:{i+1}", "source_type": "pdf"})
                return docs
            except Exception:
                # Fallback: store as raw bytes note
                return [{"doc_title": filename, "text": "(PDF解析失败：请安装 pypdf 或替换解析器)", "locator": None, "source_type": "pdf"}]

        if name.endswith(".pptx"):
            try:
                from pptx import Presentation  # type: ignore
                prs = Presentation(io.BytesIO(data))
                docs = []
                for idx, slide in enumerate(prs.slides):
                    parts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            parts.append(shape.text)
                    docs.append({"doc_title": filename, "text": "\n".join(parts), "locator": f"slide:{idx+1}", "source_type": "pptx"})
                return docs
            except Exception:
                return [{"doc_title": filename, "text": "(PPTX解析失败：请安装 python-pptx 或替换解析器)", "locator": None, "source_type": "pptx"}]

        # Fallback
        text = data.decode("utf-8", errors="ignore")
        return [{"doc_title": filename, "text": text, "locator": None, "source_type": "unknown"}]
